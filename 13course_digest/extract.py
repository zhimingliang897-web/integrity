"""
extract.py - 辅助材料文本提取模块

纯本地脚本，无需调用任何 API。
支持从 PDF 和 PPTX 文件中提取文字内容，供 LLM 分析时作为上下文参考。
"""

from pathlib import Path


def extract_pdf(pdf_path: str) -> str:
    """
    提取 PDF 文件的全部文字内容。

    Args:
        pdf_path: PDF 文件路径

    Returns:
        str: 提取的文字内容；文件不存在或解析失败时返回空字符串
    """
    if not Path(pdf_path).exists():
        print(f"[extract] PDF 文件不存在，跳过: {pdf_path}")
        return ""

    try:
        import pdfplumber
        texts: list[str] = []
        with pdfplumber.open(pdf_path) as pdf:
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages, 1):
                if i % 10 == 0 or i == total_pages:
                    print(f"[extract] 正在提取 PDF: {pdf_path} (第 {i}/{total_pages} 页)")
                text = page.extract_text()
                if text:
                    texts.append(f"[第{i}页]\n{text.strip()}")
        result = "\n\n".join(texts)
        print(f"[extract] PDF 提取完成: {pdf_path}（{len(pdf.pages)} 页）")
        return result
    except Exception as e:
        print(f"[extract] PDF 解析失败: {e}")
        return ""


def extract_pptx(pptx_path: str) -> str:
    """
    提取 PPTX 文件每张幻灯片的标题和文字内容。

    Args:
        pptx_path: PPTX 文件路径

    Returns:
        str: 按幻灯片编号排列的文字内容；文件不存在时返回空字符串
    """
    if not Path(pptx_path).exists():
        print(f"[extract] PPTX 文件不存在，跳过: {pptx_path}")
        return ""

    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        slides: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            lines: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
            if lines:
                slides.append(f"[幻灯片 {i}]\n" + "\n".join(lines))

        result = "\n\n".join(slides)
        print(f"[extract] PPTX 提取完成: {pptx_path}（{len(prs.slides)} 张）")
        return result
    except Exception as e:
        print(f"[extract] PPTX 解析失败: {e}")
        return ""


def extract_material(path: str) -> str:
    """
    根据文件扩展名自动选择提取方式（PDF 或 PPTX）。

    Args:
        path: 文件路径（.pdf / .pptx / .ppt）

    Returns:
        str: 提取的文字内容；不支持的格式返回空字符串
    """
    suffix = Path(path).suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    elif suffix in (".pptx", ".ppt"):
        return extract_pptx(path)
    else:
        print(f"[extract] 不支持的文件格式: {suffix}")
        return ""
